"""PowerPoint export with a clean academic meeting theme."""

from __future__ import annotations

import re
from datetime import date
from pathlib import Path
from typing import Any

from src.agent.state import AgentState
from src.utils.content_cleaner import clean_items


class PPTXWriter:
    """Generate a themed PPTX deck from structured report sections."""

    def __init__(self, output_dir: Path, template_path: Path | None = None) -> None:
        self.output_dir = output_dir
        self.template = self._load_template(template_path)
        self.last_warning: str | None = None

    def write_deck(self, state: AgentState, filename: str = "group_meeting_slides.pptx") -> str | None:
        """Create a PPTX file from AgentState report sections."""
        try:
            from pptx import Presentation  # type: ignore
            from pptx.dml.color import RGBColor  # type: ignore
            from pptx.enum.shapes import MSO_SHAPE  # type: ignore
            from pptx.enum.text import PP_ALIGN  # type: ignore
            from pptx.util import Inches, Pt  # type: ignore
        except ImportError:
            self.last_warning = "PPTX export skipped because python-pptx is unavailable."
            return None

        prs = Presentation()
        prs.slide_width = Inches(13.333)
        prs.slide_height = Inches(7.5)

        # Drop the default blank slide layouts into a compact namespace for helpers.
        self.Presentation = Presentation
        self.RGBColor = RGBColor
        self.MSO_SHAPE = MSO_SHAPE
        self.PP_ALIGN = PP_ALIGN
        self.Inches = Inches
        self.Pt = Pt

        slides = self._build_slide_data(state)
        for index, slide_data in enumerate(slides, start=1):
            if index == 1:
                self._add_title_slide(prs, slide_data, index, len(slides))
            elif slide_data["kind"] == "metrics":
                self._add_metrics_slide(prs, slide_data, index, len(slides))
            elif slide_data["kind"] == "paper":
                self._add_paper_slide(prs, slide_data, index, len(slides))
            elif slide_data["kind"] == "analysis":
                self._add_analysis_slide(prs, slide_data, index, len(slides))
            elif slide_data["kind"] == "plan":
                self._add_plan_slide(prs, slide_data, index, len(slides))
            elif slide_data["kind"] == "qa":
                self._add_qa_slide(prs, slide_data, index, len(slides))
            else:
                self._add_bullet_slide(prs, slide_data, index, len(slides))

        output_path = self.output_dir / filename
        try:
            prs.save(output_path)
        except PermissionError:
            self.last_warning = f"PPTX export skipped because {output_path} is locked by another application."
            return None
        return filename

    def write_from_outline(self, outline_path: Path, filename: str = "group_meeting_slides.pptx") -> str | None:
        """Backward-compatible entrypoint for older callers."""
        return None if not outline_path.exists() else None

    def _build_slide_data(self, state: AgentState) -> list[dict[str, Any]]:
        """Build eight slide payloads from report sections."""
        sections = state.report_sections
        plan = self._plan(sections)
        discussion = self._list(sections, "discussion_items", 5)
        return [
            {
                "kind": "title",
                "title": self.template["title"],
                "subtitle": self.template["subtitle"],
                "meta": f"{self.template['author']} | {date.today().isoformat()}",
            },
            {
                "kind": "bullets",
                "title": "Research Context | 研究背景",
                "bullets": self._dense_items(
                    self._list(sections, "inferred_work_items", 5),
                    [
                        "本周汇报重点放在文献启发、实验指标变化和下一步验证。",
                        "当前结论需要结合更多场景和失败案例继续确认。",
                        "展示主线建议围绕研究问题、证据和风险展开。",
                    ],
                    limit=5,
                ),
            },
            {
                "kind": "paper",
                "title": "Paper Takeaways | 文献启发",
                "cards": self._dicts(sections, "paper_cards")[:3],
            },
            {
                "kind": "metrics",
                "title": "Experiment Evidence | 实验证据",
                "metrics": self._dicts(sections, "metric_cards")[:4],
                "bullets": self._dense_items(
                    self._list(sections, "experiment_items", 5),
                    [
                        "指标不足时，可结合结果图或失败案例解释实验现象。",
                        "实验结论应同时关注图像质量指标和下游检测指标。",
                        "下一轮实验建议围绕关键假设补齐 before/after 证据。",
                    ],
                    limit=5,
                ),
            },
            {
                "kind": "analysis",
                "title": "Key Insights | 关键洞察",
                "cards": self._dicts(sections, "analysis_cards")[:3],
            },
            {
                "kind": "bullets",
                "title": "Issues & Risks | 当前问题与风险",
                "bullets": self._dense_items(
                    self._list(sections, "issue_items", 5),
                    [
                        "文献启发和当前实验之间的对应关系仍需人工确认。",
                        "指标提升是否稳定仍需更多场景验证。",
                        "失败案例不足时，风险分析需要保持保守。",
                    ],
                    limit=5,
                ),
            },
            {
                "kind": "plan",
                "title": "Next Actions | 下一步行动",
                "plan": plan,
            },
            {
                "kind": "bullets",
                "title": "Discussion | 需要导师建议",
                "bullets": self._dense_items(
                    discussion,
                    [
                        "下周是否优先补充低光照和运动模糊场景？",
                        "是否建议将检测连续性作为核心汇报指标之一？",
                        "是否需要增加失败案例可视化来支撑指标解释？",
                    ],
                    limit=5,
                ),
            },
        ]

    def _add_title_slide(self, prs: Any, data: dict[str, Any], page: int, total: int) -> None:
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        self._set_background(slide)
        self._add_accent_bar(slide)
        self._add_textbox(slide, data["title"], 0.85, 1.55, 11.6, 0.8, 46, bold=True, color="title")
        self._add_textbox(slide, data["subtitle"], 0.9, 2.55, 10.8, 0.5, 23, color="body")
        self._add_textbox(slide, data["meta"], 0.9, 3.25, 10.8, 0.4, 16, color="muted")
        self._add_footer(slide, page, total)

    def _add_bullet_slide(self, prs: Any, data: dict[str, Any], page: int, total: int) -> None:
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        self._prepare_content_slide(slide, data["title"], page, total)
        self._add_card(slide, 0.75, 1.45, 11.85, 4.85)
        self._add_bullets(slide, data.get("bullets", []), 1.1, 1.8, 11.1, 4.2, font_size=19)

    def _add_paper_slide(self, prs: Any, data: dict[str, Any], page: int, total: int) -> None:
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        self._prepare_content_slide(slide, data["title"], page, total)
        cards = data.get("cards") or []
        y = 1.35
        for index, card in enumerate(cards[:3], start=1):
            self._add_card(slide, 0.75, y, 11.85, 1.55)
            title = f"{index}. {card.get('paper', '文献')}"
            details = [
                f"论文：{card.get('paper', '')}",
                f"方法：{card.get('method', '')}",
                f"启发：{card.get('insight', '')}",
                f"局限：{card.get('limitation', '')}",
            ]
            self._add_textbox(slide, title, 1.05, y + 0.12, 11.2, 0.3, 17, bold=True, color="title")
            self._add_bullets(slide, details[:4], 1.1, y + 0.48, 11.0, 0.95, font_size=12)
            y += 1.72

    def _add_metrics_slide(self, prs: Any, data: dict[str, Any], page: int, total: int) -> None:
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        self._prepare_content_slide(slide, data["title"], page, total)
        metrics = data.get("metrics") or []
        card_w = 2.85
        for idx, metric in enumerate(metrics[:4]):
            x = 0.75 + idx * 3.03
            self._add_card(slide, x, 1.35, card_w, 1.75, fill="soft")
            self._add_textbox(slide, str(metric.get("name")), x + 0.18, 1.55, card_w - 0.35, 0.3, 18, bold=True, color="title")
            before = self._fmt(metric.get("before"))
            after = self._fmt(metric.get("after"))
            delta = self._fmt(metric.get("delta"), signed=True)
            unit = str(metric.get("unit", ""))
            self._add_textbox(slide, f"{before} -> {after}", x + 0.18, 2.0, card_w - 0.35, 0.35, 20, bold=True, color="primary")
            self._add_textbox(slide, f"提升 {delta}{unit}", x + 0.18, 2.45, card_w - 0.35, 0.28, 14, color="muted")
        self._add_card(slide, 0.75, 3.45, 11.85, 2.0)
        self._add_bullets(slide, data.get("bullets", []), 1.1, 3.75, 11.0, 1.45, font_size=17)

    def _add_analysis_slide(self, prs: Any, data: dict[str, Any], page: int, total: int) -> None:
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        self._prepare_content_slide(slide, data["title"], page, total)
        headers = ["发现", "解释", "证据"]
        col_w = 3.75
        for col, header in enumerate(headers):
            x = 0.75 + col * 3.95
            self._add_textbox(slide, header, x, 1.32, col_w, 0.35, 18, bold=True, color="primary")
        y = 1.78
        for card in data.get("cards", [])[:3]:
            self._add_card(slide, 0.75, y, 11.85, 1.15)
            values = [card.get("finding", ""), card.get("explanation", ""), card.get("evidence", "")]
            for col, value in enumerate(values):
                self._add_textbox(slide, str(value), 0.95 + col * 3.95, y + 0.18, 3.35, 0.75, 13, color="body")
            y += 1.32

    def _add_plan_slide(self, prs: Any, data: dict[str, Any], page: int, total: int) -> None:
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        self._prepare_content_slide(slide, data["title"], page, total)
        plan = data.get("plan", {})
        columns = [("Must do", plan.get("must", [])), ("Should do", plan.get("should", [])), ("Could do", plan.get("could", []))]
        for idx, (label, items) in enumerate(columns):
            x = 0.75 + idx * 4.05
            self._add_card(slide, x, 1.35, 3.75, 4.95, fill="soft" if idx == 0 else "white")
            self._add_textbox(slide, label, x + 0.25, 1.6, 3.25, 0.35, 19, bold=True, color="title")
            self._add_bullets(slide, items[:4], x + 0.25, 2.1, 3.25, 3.6, font_size=14)

    def _add_qa_slide(self, prs: Any, data: dict[str, Any], page: int, total: int) -> None:
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        self._prepare_content_slide(slide, data["title"], page, total)
        questions = data.get("questions", [])
        left = [f"{q} | {a}" for q, a in questions[:4]]
        right = [f"{q} | {a}" for q, a in questions[4:8]]
        self._add_card(slide, 0.75, 1.35, 5.75, 4.95)
        self._add_card(slide, 6.85, 1.35, 5.75, 4.95)
        self._add_bullets(slide, left, 1.05, 1.7, 5.2, 4.25, font_size=13)
        self._add_bullets(slide, right, 7.15, 1.7, 5.2, 4.25, font_size=13)

    def _prepare_content_slide(self, slide: Any, title: str, page: int, total: int) -> None:
        self._set_background(slide)
        self._add_accent_bar(slide)
        self._add_textbox(slide, title, 0.75, 0.45, 11.8, 0.55, 30, bold=True, color="title")
        self._add_footer(slide, page, total)

    def _set_background(self, slide: Any) -> None:
        fill = slide.background.fill
        fill.solid()
        fill.fore_color.rgb = self._rgb("background")

    def _add_accent_bar(self, slide: Any) -> None:
        shape = slide.shapes.add_shape(self.MSO_SHAPE.RECTANGLE, self.Inches(0), self.Inches(0), self.Inches(13.333), self.Inches(0.12))
        shape.fill.solid()
        shape.fill.fore_color.rgb = self._rgb("primary")
        shape.line.fill.background()

    def _add_footer(self, slide: Any, page: int, total: int) -> None:
        if self.template.get("show_footer", True):
            self._add_textbox(slide, self.template.get("footer_text", "Research Meeting"), 0.75, 7.05, 5.0, 0.25, 10, color="muted")
        if self.template.get("show_page_number", True):
            self._add_textbox(slide, f"{page}/{total}", 11.8, 7.05, 0.75, 0.25, 10, color="muted", align="right")

    def _add_card(self, slide: Any, x: float, y: float, width: float, height: float, fill: str = "white") -> Any:
        shape = slide.shapes.add_shape(
            self.MSO_SHAPE.ROUNDED_RECTANGLE,
            self.Inches(x),
            self.Inches(y),
            self.Inches(width),
            self.Inches(height),
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = self._rgb("soft" if fill == "soft" else "card")
        shape.line.color.rgb = self._rgb("line")
        shape.line.width = self.Pt(0.75)
        return shape

    def _add_bullets(
        self, slide: Any, items: list[str], x: float, y: float, width: float, height: float, font_size: int = 16
    ) -> None:
        textbox = slide.shapes.add_textbox(self.Inches(x), self.Inches(y), self.Inches(width), self.Inches(height))
        frame = textbox.text_frame
        frame.clear()
        frame.word_wrap = True
        safe_items = self._visible_items([str(item) for item in items], limit=8) or ["暂无明确记录。"]
        for index, item in enumerate(safe_items):
            paragraph = frame.paragraphs[0] if index == 0 else frame.add_paragraph()
            paragraph.text = item
            paragraph.level = 0
            paragraph.font.name = self.template["primary_font"]
            paragraph.font.size = self.Pt(font_size)
            paragraph.font.color.rgb = self._rgb("body")
            paragraph.space_after = self.Pt(5)

    def _add_textbox(
        self,
        slide: Any,
        text: str,
        x: float,
        y: float,
        width: float,
        height: float,
        font_size: int,
        bold: bool = False,
        color: str = "body",
        align: str = "left",
    ) -> None:
        textbox = slide.shapes.add_textbox(self.Inches(x), self.Inches(y), self.Inches(width), self.Inches(height))
        frame = textbox.text_frame
        frame.clear()
        frame.word_wrap = True
        paragraph = frame.paragraphs[0]
        paragraph.text = str(text)
        paragraph.font.name = self.template["primary_font"]
        paragraph.font.size = self.Pt(font_size)
        paragraph.font.bold = bold
        paragraph.font.color.rgb = self._rgb(color)
        paragraph.alignment = self.PP_ALIGN.RIGHT if align == "right" else self.PP_ALIGN.LEFT

    def _load_template(self, template_path: Path | None) -> dict[str, Any]:
        """Load a YAML theme file with safe defaults."""
        defaults: dict[str, Any] = {
            "title": "本周组会汇报",
            "subtitle": "视频增强 / 目标检测 / 实验评估",
            "author": "Graduate Research Workflow",
            "theme_name": "academic_clean",
            "footer_text": "Research Meeting",
            "primary_font": "Microsoft YaHei",
            "slide_size": "16:9",
            "show_page_number": True,
            "show_footer": True,
        }
        if not template_path or not template_path.exists():
            return defaults
        try:
            import yaml  # type: ignore

            loaded = yaml.safe_load(template_path.read_text(encoding="utf-8")) or {}
            defaults.update(loaded)
        except Exception:
            return defaults
        return defaults

    def _rgb(self, name: str) -> Any:
        colors = {
            "background": (248, 250, 252),
            "card": (255, 255, 255),
            "soft": (235, 244, 255),
            "line": (203, 213, 225),
            "title": (15, 23, 42),
            "body": (51, 65, 85),
            "muted": (100, 116, 139),
            "primary": (37, 99, 235),
        }
        return self.RGBColor(*colors.get(name, colors["body"]))

    def _list(self, sections: dict[str, object], key: str, limit: int) -> list[str]:
        values = sections.get(key, [])
        return self._visible_items([str(item) for item in values], limit=limit) if isinstance(values, list) else []

    def _dense_items(self, items: list[str], fallbacks: list[str], limit: int = 5) -> list[str]:
        """Keep slides useful even when a source file is an empty template."""
        merged = clean_items(items, limit=limit)
        for item in fallbacks:
            if len(merged) >= min(3, limit):
                break
            merged.append(item)
        return self._visible_items(merged, limit=limit)

    def _dicts(self, sections: dict[str, object], key: str) -> list[dict[str, Any]]:
        values = sections.get(key, [])
        return [item for item in values if isinstance(item, dict)] if isinstance(values, list) else []

    def _plan(self, sections: dict[str, object]) -> dict[str, list[str]]:
        raw = sections.get("prioritized_plan", {})
        if not isinstance(raw, dict):
            raw = {}
        return {
            "must": self._visible_items([str(item) for item in raw.get("must", [])], limit=4),
            "should": self._visible_items([str(item) for item in raw.get("should", [])], limit=4),
            "could": self._visible_items([str(item) for item in raw.get("could", [])], limit=4),
        }

    def _questions(self, sections: dict[str, object]) -> list[tuple[str, str]]:
        raw = sections.get("questions", [])
        result: list[tuple[str, str]] = []
        if isinstance(raw, list):
            for item in raw:
                if isinstance(item, (list, tuple)) and len(item) >= 2:
                    result.append((str(item[0]), str(item[1])))
        return result

    def _fmt(self, value: object, signed: bool = False) -> str:
        try:
            number = float(value)
        except (TypeError, ValueError):
            return "N/A"
        prefix = "+" if signed and number > 0 else ""
        return f"{prefix}{number:.3f}".rstrip("0").rstrip(".")

    def _visible_items(self, items: list[str], limit: int | None = None) -> list[str]:
        patterns = [
            r"检测到\s*\d+\s*篇\s*PDF",
            r"空模板",
            r"占位内容已过滤",
            r"progress\.md\s*为空",
            r"agent_state\.json\s*记录",
            r"use_paper_notes_as_sources",
            r"source_type",
            r"本工具",
            r"本\s*Agent",
            r"运行命令",
        ]
        visible = []
        for item in items:
            text = str(item).strip()
            if any(re.search(pattern, text, flags=re.IGNORECASE) for pattern in patterns):
                continue
            visible.append(text)
        return clean_items(visible, limit=limit, min_chars=4)
