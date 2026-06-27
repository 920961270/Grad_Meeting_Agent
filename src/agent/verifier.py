"""Output verification and quality checks."""

from __future__ import annotations

import re
from pathlib import Path


class Verifier:
    """Check generated files for required content and quality gates."""

    INTERNAL_PATTERNS = [
        r"检测到\s*\d+\s*篇\s*PDF",
        r"空模板",
        r"占位内容已过滤",
        r"progress\.md\s*为空",
        r"agent_state\.json\s*记录",
        r"use_paper_notes_as_sources",
        r"source_type",
        r"本工具",
        r"本\s*Agent",
        r"运行命令",
        r"请补充\s*progress",
        r"请补充真实进展",
        r"需要补充真实实验进展",
        r"请先填写\s*progress",
    ]

    def __init__(
        self,
        output_dir: Path,
        input_dir: Path | None = None,
        use_paper_notes_as_sources: bool = False,
        detected_pdf_papers: list[str] | None = None,
        with_slides: bool = False,
        export_details: bool = False,
    ) -> None:
        self.output_dir = output_dir
        self.input_dir = input_dir
        self.use_paper_notes_as_sources = use_paper_notes_as_sources
        self.detected_pdf_papers = detected_pdf_papers or []
        self.with_slides = with_slides
        self.export_details = export_details

    def verify(self) -> dict[str, object]:
        """Run structural checks and quality checks."""
        checks = {
            "meeting_pack": self._check_meeting_pack(),
            "default_output_contract": self._check_output_contract(),
        }
        quality_checks = {
            "meeting_pack_clean": self._check_visible_clean(["meeting_pack.md"]),
            "meeting_pack_has_supervisor_prep": self._check_meeting_pack_supervisor_prep(),
            "no_template_placeholders": self._check_no_template_placeholders(),
        }

        if self.with_slides:
            checks.update(
                {
                    "ppt_outline": self._check_ppt_outline(),
                    "pptx": self._check_pptx(),
                }
            )
            quality_checks.update(
                {
                    "ppt_outline_density": self._check_ppt_outline_density(),
                    "pptx_no_supervisor_qa": self._check_pptx_no_supervisor_qa(),
                    "pptx_no_internal_terms": self._check_pptx_no_internal_terms(),
                    "pptx_final_slide_discussion": self._check_pptx_final_slide_discussion(),
                }
            )

        if self.export_details:
            checks.update(
                {
                    "details_research_brief": self._check_non_empty("details/research_brief.md"),
                    "details_paper_cards": self._check_non_empty("details/paper_cards.md"),
                    "details_experiment_insights": self._check_non_empty("details/experiment_insights.md"),
                    "details_meeting_agenda": self._check_non_empty("details/meeting_agenda.md"),
                    "details_weekly_report": self._check_weekly_report("details/weekly_report.md"),
                    "details_oral_script": self._check_non_empty("details/oral_script.md"),
                    "details_supervisor_qa": self._check_supervisor_qa("details/supervisor_qa.md"),
                    "details_next_week_plan": self._check_next_week_plan("details/next_week_plan.md"),
                }
            )
            quality_checks["details_clean"] = self._check_visible_clean(
                [
                    "details/research_brief.md",
                    "details/paper_cards.md",
                    "details/experiment_insights.md",
                    "details/meeting_agenda.md",
                    "details/weekly_report.md",
                    "details/oral_script.md",
                    "details/next_week_plan.md",
                ]
            )

        passed = all(item["passed"] for item in checks.values()) and all(
            item["passed"] for item in quality_checks.values()
        )
        return {"passed": passed, "checks": checks, "quality_checks": quality_checks}

    def _read(self, filename: str) -> str:
        path = self.output_dir / filename
        if not path.exists():
            return ""
        return path.read_text(encoding="utf-8")

    def _check_meeting_pack(self) -> dict[str, object]:
        text = self._read("meeting_pack.md")
        required = [
            "# Research Meeting Pack",
            "## 1. 一页概览",
            "## 2. 本周可以汇报什么",
            "## 3. 文献卡片",
            "## 4. 实验结果洞察",
            "## 5. 图像与结果材料",
            "## 6. 当前问题与风险",
            "## 7. 建议 PPT 结构",
            "## 8. 组会讨论问题",
            "## 9. 个人准备：导师可能追问",
        ]
        missing = [item for item in required if item not in text]
        return {"passed": bool(text) and not missing, "missing": missing}

    def _check_output_contract(self) -> dict[str, object]:
        root_files = [
            path.name
            for path in self.output_dir.iterdir()
            if path.is_file() and path.name != "agent_state.json"
        ]
        expected = {"meeting_pack.md"}
        if self.with_slides:
            expected.update({"ppt_outline.md", "group_meeting_slides.pptx"})
        unexpected = sorted(name for name in root_files if name not in expected)
        details_exists = (self.output_dir / "details").exists()
        details_ok = self.export_details or not details_exists
        return {
            "passed": not unexpected and details_ok,
            "root_files": sorted(root_files),
            "unexpected": unexpected,
            "details_exists": details_exists,
        }

    def _check_visible_clean(self, filenames: list[str]) -> dict[str, object]:
        combined = "\n".join(self._read(filename) for filename in filenames)
        findings = self._find_patterns(combined, self.INTERNAL_PATTERNS)
        return {"passed": not findings, "findings": findings}

    def _check_no_template_placeholders(self) -> dict[str, object]:
        files = ["meeting_pack.md"]
        if self.with_slides:
            files.append("ppt_outline.md")
        if self.export_details:
            files.extend(
                [
                    "details/research_brief.md",
                    "details/paper_cards.md",
                    "details/experiment_insights.md",
                    "details/meeting_agenda.md",
                    "details/weekly_report.md",
                    "details/oral_script.md",
                    "details/next_week_plan.md",
                ]
            )
        combined = "\n".join(self._read(filename) for filename in files)
        bad_patterns = [
            r"请在这里填写",
            r"关注点：\s*$",
            r"方法启发：\s*$",
            r"和当前实验关系：\s*$",
            r"疑问：\s*$",
            r"^\s*[-*]\s*$",
        ]
        findings = self._find_patterns(combined, bad_patterns)
        return {"passed": not findings, "findings": findings}

    def _check_meeting_pack_supervisor_prep(self) -> dict[str, object]:
        text = self._read("meeting_pack.md")
        questions = re.findall(r"^###\s+Q\d+", text, flags=re.MULTILINE)
        return {"passed": len(questions) >= 6, "question_count": len(questions)}

    def _check_weekly_report(self, filename: str) -> dict[str, object]:
        text = self._read(filename)
        required = ["本周工作概览", "文献阅读", "实验/项目进展", "问题与分析", "下周计划"]
        missing = [item for item in required if item not in text]
        return {"passed": bool(text) and not missing, "missing": missing}

    def _check_ppt_outline(self) -> dict[str, object]:
        text = self._read("ppt_outline.md")
        slides = re.findall(r"^##\s+Slide\s+\d+", text, flags=re.MULTILINE)
        return {"passed": len(slides) >= 8, "slide_count": len(slides)}

    def _check_ppt_outline_density(self) -> dict[str, object]:
        text = self._read("ppt_outline.md")
        sections = re.split(r"^##\s+Slide\s+\d+:", text, flags=re.MULTILINE)[1:]
        sparse: list[int] = []
        for index, section in enumerate(sections, start=1):
            bullets = re.findall(r"^\s*[-*]\s+", section, flags=re.MULTILINE)
            if index != 1 and len(bullets) < 3:
                sparse.append(index)
        return {"passed": len(sections) >= 8 and not sparse, "slide_count": len(sections), "sparse_slides": sparse}

    def _check_pptx(self) -> dict[str, object]:
        slide_texts = self._pptx_slide_texts()
        return {"passed": len(slide_texts) >= 8, "slide_count": len(slide_texts)}

    def _check_pptx_no_supervisor_qa(self) -> dict[str, object]:
        text = "\n".join(self._pptx_slide_texts())
        bad_terms = ["Supervisor Q&A", "导师可能追问", "个人准备"]
        findings = [term for term in bad_terms if term in text]
        return {"passed": not findings, "findings": findings}

    def _check_pptx_no_internal_terms(self) -> dict[str, object]:
        text = "\n".join(self._pptx_slide_texts())
        findings = self._find_patterns(text, self.INTERNAL_PATTERNS)
        return {"passed": not findings, "findings": findings}

    def _check_pptx_final_slide_discussion(self) -> dict[str, object]:
        slides = self._pptx_slide_texts()
        if not slides:
            return {"passed": False, "last_slide": ""}
        last = slides[-1]
        return {"passed": "Discussion" in last, "last_slide": last[:300]}

    def _check_supervisor_qa(self, filename: str) -> dict[str, object]:
        text = self._read(filename)
        questions = re.findall(r"^###\s+Q\d+", text, flags=re.MULTILINE)
        return {"passed": len(questions) >= 8, "question_count": len(questions)}

    def _check_non_empty(self, filename: str) -> dict[str, object]:
        text = self._read(filename).strip()
        return {"passed": bool(text), "chars": len(text)}

    def _check_next_week_plan(self, filename: str) -> dict[str, object]:
        text = self._read(filename)
        required = ["Must do", "Should do", "Could do"]
        missing = [item for item in required if item not in text]
        bullet_count = len(re.findall(r"^- ", text, flags=re.MULTILINE))
        return {"passed": not missing and bullet_count >= 6, "missing": missing, "tasks": bullet_count}

    def _pptx_slide_texts(self) -> list[str]:
        path = self.output_dir / "group_meeting_slides.pptx"
        if not path.exists():
            return []
        try:
            from pptx import Presentation  # type: ignore

            deck = Presentation(str(path))
        except Exception:
            return []
        slides: list[str] = []
        for slide in deck.slides:
            parts = []
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    parts.append(shape.text.strip())
            slides.append("\n".join(parts))
        return slides

    def _find_patterns(self, text: str, patterns: list[str]) -> list[str]:
        findings: list[str] = []
        for pattern in patterns:
            if re.search(pattern, text, flags=re.IGNORECASE | re.MULTILINE):
                findings.append(pattern)
        return findings
